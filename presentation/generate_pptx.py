#!/usr/bin/env python3
"""
Generate kage_5min.pptx from slides/<NN>_<name>/{slide.md, pointers.md}.

slide.md syntax:
    # Title                               -> slide title
    <!-- LAYOUT: architecture -->         -> render the architecture diagram
                                             instead of bullets (native shapes,
                                             every box/arrow is editable in PPT)
    ## Subtitle (optional)
    - bullet
    - bullet
    ```python
    code...
    ```
    > Quote line

pointers.md becomes the slide's speaker notes verbatim.

Run:
    pip install python-pptx
    python presentation/generate_pptx.py
"""
from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent
SLIDES_DIR = ROOT / "slides"
OUT_PPTX = ROOT / "kage_5min.pptx"

# Light / white theme palette
BG       = RGBColor(0xFF, 0xFF, 0xFF)   # slide background
TEXT     = RGBColor(0x1F, 0x29, 0x37)   # primary text
MUTED    = RGBColor(0x6B, 0x72, 0x80)   # secondary text
ACCENT   = RGBColor(0x25, 0x63, 0xEB)   # blue accent
SUCCESS  = RGBColor(0x05, 0x96, 0x69)   # green
CODEBG   = RGBColor(0xF6, 0xF8, 0xFA)   # very light grey for code panels
PANEL    = RGBColor(0xF3, 0xF4, 0xF6)   # card fill
BORDER   = RGBColor(0xD1, 0xD5, 0xDB)   # card border
GOLD     = RGBColor(0xB4, 0x53, 0x09)   # for inline code text

LAYER_FILLS = {
    "User Code":   RGBColor(0xE0, 0xF2, 0xFE),  # sky-100
    "KAGE API":    RGBColor(0xDB, 0xEA, 0xFE),  # blue-100
    "Adapters":    RGBColor(0xDC, 0xFC, 0xE7),  # green-100
    "Core Engine": RGBColor(0xEE, 0xE9, 0xFE),  # violet-100
    "Storage":     RGBColor(0xFE, 0xF3, 0xC7),  # amber-100
    "Analytics":   RGBColor(0xFE, 0xE2, 0xE2),  # red-100
}


# --------------------------------------------------------------------- parser

def _is_table_separator(line: str) -> bool:
    s = line.strip()
    if not (s.startswith("|") and s.endswith("|")):
        return False
    cells = [c.strip() for c in s.strip("|").split("|")]
    return all(re.fullmatch(r":?-{3,}:?", c) for c in cells)


def parse_slide(path: Path) -> dict:
    text = path.read_text()
    lines = text.splitlines()

    title = ""
    subtitle = ""
    bullets: list[str] = []
    code_blocks: list[str] = []
    quote = ""
    layout = "default"
    table_rows: list[list[str]] = []

    i = 0
    in_code = False
    code_buf: list[str] = []

    while i < len(lines):
        line = lines[i]
        if line.startswith("```"):
            if in_code:
                code_blocks.append("\n".join(code_buf))
                code_buf = []
                in_code = False
            else:
                in_code = True
            i += 1
            continue
        if in_code:
            code_buf.append(line)
            i += 1
            continue

        m = re.match(r"<!--\s*LAYOUT:\s*(\w+)\s*-->", line.strip())
        if m:
            layout = m.group(1)
            i += 1
            continue

        stripped = line.strip()
        if stripped.startswith("|") and stripped.endswith("|"):
            if _is_table_separator(line):
                i += 1
                continue
            cells = [c.strip() for c in stripped.strip("|").split("|")]
            table_rows.append(cells)
            i += 1
            continue

        if line.startswith("# ") and not title:
            title = line[2:].strip()
        elif line.startswith("## ") and not subtitle:
            subtitle = line[3:].strip()
        elif line.startswith("- "):
            bullets.append(line[2:].strip())
        elif line.startswith("> "):
            quote = line[2:].strip()
        i += 1

    return {
        "title": title,
        "subtitle": subtitle,
        "bullets": bullets,
        "code": code_blocks[0] if code_blocks else None,
        "quote": quote,
        "layout": layout,
        "table": table_rows,
    }


# --------------------------------------------------------------------- helpers

def set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, *,
                size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT,
                font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return tb


def add_bullets(slide, left, top, width, height, bullets, *,
                size=20, color=TEXT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True

    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        parts = re.split(r"(\*\*[^*]+\*\*|`[^`]+`)", bullet)
        # leading bullet
        lead = p.add_run()
        lead.text = "•  "
        lead.font.size = Pt(size)
        lead.font.color.rgb = ACCENT
        lead.font.name = "Calibri"
        for part in parts:
            run = p.add_run()
            if part.startswith("**") and part.endswith("**"):
                run.text = part[2:-2]
                run.font.bold = True
                run.font.color.rgb = TEXT
            elif part.startswith("`") and part.endswith("`"):
                run.text = part[1:-1]
                run.font.name = "Menlo"
                run.font.color.rgb = GOLD
            else:
                run.text = part
                run.font.color.rgb = color
            run.font.size = Pt(size)
            if not run.font.name:
                run.font.name = "Calibri"


def add_code_block(slide, left, top, width, height, code):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CODEBG
    shape.line.color.rgb = BORDER
    shape.shadow.inherit = False

    tb = slide.shapes.add_textbox(
        left + Inches(0.18), top + Inches(0.12),
        width - Inches(0.36), height - Inches(0.24),
    )
    tf = tb.text_frame
    tf.word_wrap = False
    for idx, line in enumerate(code.splitlines() or [""]):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line if line else " "
        run.font.name = "Menlo"
        run.font.size = Pt(12)
        run.font.color.rgb = TEXT


def add_notes(slide, notes: str) -> None:
    slide.notes_slide.notes_text_frame.text = notes


# --------------------------------------------------------------------- shapes

def _layer_card(slide, left, top, width, height, title, body, fill_rgb):
    """A rounded card with a title and a small body line. Fully editable."""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = fill_rgb
    box.line.color.rgb = BORDER
    box.line.width = Pt(0.75)
    box.shadow.inherit = False
    box.name = f"layer__{title.replace(' ', '_')}"

    tf = box.text_frame
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    tf.word_wrap = True

    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.LEFT
    r1 = p1.add_run()
    r1.text = title
    r1.font.size = Pt(14)
    r1.font.bold = True
    r1.font.color.rgb = TEXT
    r1.font.name = "Calibri"

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    r2.text = body
    r2.font.size = Pt(11)
    r2.font.color.rgb = MUTED
    r2.font.name = "Calibri"


def _down_arrow(slide, left, top, length=Inches(0.28)):
    """A short down-arrow connector between layers. Editable."""
    arrow = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        left, top, left, top + length,
    )
    line = arrow.line
    line.color.rgb = MUTED
    line.width = Pt(2.25)
    arrow.name = "arrow_down"
    # end arrowhead via XML — python-pptx exposes only basic line props
    from pptx.oxml.ns import qn
    ln = arrow.line._get_or_add_ln()
    tail = ln.find(qn("a:tailEnd"))
    if tail is None:
        from lxml import etree
        tail = etree.SubElement(ln, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("h", "med")


def _style_table_cell(cell, text, *, bold=False, size=10, color=TEXT,
                      fill=None, align=PP_ALIGN.LEFT):
    cell.text = ""
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    tf = cell.text_frame
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def build_comparison_slide(prs, slide_data, notes, idx, total):
    """Render a markdown table as a native, fully-editable PowerPoint table."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, BG)
    add_accent_stripe(slide)
    add_page_number(slide, idx, total)

    add_textbox(slide, Inches(0.5), Inches(0.35), Inches(12.5), Inches(0.7),
                slide_data["title"], size=28, bold=True, color=TEXT)
    if slide_data["subtitle"]:
        add_textbox(slide, Inches(0.5), Inches(1.05), Inches(12.5), Inches(0.5),
                    slide_data["subtitle"], size=16, color=ACCENT)

    rows = slide_data["table"]
    if not rows:
        add_notes(slide, notes)
        return

    n_cols = max(len(r) for r in rows)
    # normalise row widths
    for r in rows:
        while len(r) < n_cols:
            r.append("")

    n_rows = len(rows)
    table_left = Inches(0.5)
    table_top = Inches(1.65)
    table_w = Inches(12.33)
    # quote takes a strip at the bottom if present
    bottom_reserved = Inches(0.6) if slide_data["quote"] else Inches(0.25)
    table_h = Inches(7.5) - table_top - bottom_reserved

    shape = slide.shapes.add_table(n_rows, n_cols, table_left, table_top,
                                   table_w, table_h)
    table = shape.table
    shape.name = "kage_databricks_comparison"

    # Column widths — first column narrower (it's the label column)
    if n_cols == 3:
        table.columns[0].width = Inches(3.0)
        table.columns[1].width = Inches(4.6)
        table.columns[2].width = Inches(4.73)
    else:
        even = int(table_w / n_cols)
        for c in table.columns:
            c.width = Emu(even)

    for ri, row in enumerate(rows):
        for ci, text in enumerate(row):
            cell = table.cell(ri, ci)
            if ri == 0:
                _style_table_cell(cell, text, bold=True, size=12,
                                  color=RGBColor(0xFF, 0xFF, 0xFF),
                                  fill=ACCENT)
            else:
                fill = PANEL if ri % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
                # accent the "What KAGE adds" column subtly
                color = TEXT
                if ci == n_cols - 1:
                    color = TEXT
                _style_table_cell(cell, text, size=10, color=color, fill=fill)

    if slide_data["quote"]:
        add_textbox(slide, Inches(0.5), Inches(6.85), Inches(12.5), Inches(0.45),
                    f"“{slide_data['quote']}”",
                    size=13, color=MUTED, align=PP_ALIGN.CENTER)

    add_notes(slide, notes)


def build_architecture_slide(prs, slide_data, notes, idx, total):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, BG)

    add_accent_stripe(slide)
    add_page_number(slide, idx, total)

    # Title + subtitle
    add_textbox(slide, Inches(0.5), Inches(0.35), Inches(12.5), Inches(0.7),
                slide_data["title"], size=32, bold=True, color=TEXT)
    if slide_data["subtitle"]:
        add_textbox(slide, Inches(0.5), Inches(1.05), Inches(12.5), Inches(0.5),
                    slide_data["subtitle"], size=18, color=ACCENT)

    # Architecture stack: 6 layers, centered. Each layer + arrow editable.
    layers = [
        ("User Code",
         "PySpark · DLT · dbt · Airflow · LangChain / LangGraph · async agents"),
        ("KAGE API",
         "ETL: @pipeline / @task / @dataset    Agentic: @agent / @step / @tool / @llm_call"),
        ("Adapters",
         "@kage_dlt_table · emit_dbt_run_results · kage_task_callbacks · KageLangChainCallback"),
        ("Core Engine",
         "thread-safe lock · contextvars span stack · auto error capture · schemas"),
        ("Storage",
         "partitioned JSONL: base/{platform}/event_type=X/dt=Y/part-*.jsonl"),
        ("Analytics",
         "SLA · Medallion Health · Lineage · Cost · Agentic call tree · AI suggestions"),
    ]

    card_w = Inches(9.0)
    card_h = Inches(0.62)
    gap    = Inches(0.16)   # space for the arrow between cards
    left   = Inches(2.15)   # centered: (13.33 - 9.0) / 2 ≈ 2.165
    top    = Inches(1.75)

    for title, body in layers:
        _layer_card(slide, left, top, card_w, card_h, title, body,
                    LAYER_FILLS.get(title, PANEL))
        # arrow centered horizontally on the card, sitting just below
        arrow_x = left + Inches(4.5)
        _down_arrow(slide, arrow_x, top + card_h + Emu(20000), length=gap - Emu(40000))
        top = top + card_h + gap

    # Side caption
    add_textbox(slide, Inches(0.4), Inches(6.7), Inches(12.5), Inches(0.4),
                "Every shape on this slide is a native PowerPoint object — move, recolour, or relabel freely.",
                size=11, color=MUTED, align=PP_ALIGN.CENTER)

    add_notes(slide, notes)


# --------------------------------------------------------------------- default slide

def add_accent_stripe(slide):
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                    Inches(0.16), Inches(7.5))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = ACCENT
    stripe.line.fill.background()
    stripe.shadow.inherit = False


def add_page_number(slide, idx, total):
    label = f"{idx} / {total}"
    add_textbox(slide, Inches(12.2), Inches(7.1), Inches(1.0), Inches(0.3),
                label, size=11, color=MUTED, align=PP_ALIGN.RIGHT)


def build_default_slide(prs, slide_data, notes, idx, total):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, BG)
    add_accent_stripe(slide)
    add_page_number(slide, idx, total)

    add_textbox(slide, Inches(0.5), Inches(0.4), Inches(12.5), Inches(0.9),
                slide_data["title"], size=36, bold=True, color=TEXT)

    y = Inches(1.5)
    if slide_data["subtitle"]:
        add_textbox(slide, Inches(0.5), y, Inches(12.5), Inches(0.6),
                    slide_data["subtitle"], size=22, color=ACCENT)
        y = Inches(2.2)

    if slide_data["code"]:
        n_lines = len(slide_data["code"].splitlines())
        code_h = min(Inches(0.30) * n_lines + Inches(0.4), Inches(4.6))
        add_code_block(slide, Inches(0.5), y, Inches(12.5), code_h, slide_data["code"])
        y = y + code_h + Inches(0.2)
        if slide_data["bullets"]:
            add_bullets(slide, Inches(0.5), y, Inches(12.5),
                        Inches(7.5) - y - Inches(0.7),
                        slide_data["bullets"], size=15)
    elif slide_data["bullets"]:
        add_bullets(slide, Inches(0.5), y, Inches(12.5),
                    Inches(7.5) - y - Inches(0.8),
                    slide_data["bullets"], size=22)

    if slide_data["quote"]:
        add_textbox(slide, Inches(0.5), Inches(6.6), Inches(12.5), Inches(0.5),
                    f"“{slide_data['quote']}”",
                    size=16, color=MUTED)

    add_notes(slide, notes)


# --------------------------------------------------------------------- main

def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    folders = sorted(p for p in SLIDES_DIR.iterdir() if p.is_dir())
    folders = [f for f in folders if (f / "slide.md").exists()]
    if not folders:
        raise SystemExit(f"No slide folders under {SLIDES_DIR}")

    total = len(folders)

    for idx, folder in enumerate(folders, start=1):
        slide_md = folder / "slide.md"
        pointers_md = folder / "pointers.md"
        data = parse_slide(slide_md)
        notes = pointers_md.read_text() if pointers_md.exists() else ""

        if data["layout"] == "architecture":
            build_architecture_slide(prs, data, notes, idx, total)
            marker = "[arch]"
        elif data["layout"] == "comparison":
            build_comparison_slide(prs, data, notes, idx, total)
            marker = "[cmp] "
        else:
            build_default_slide(prs, data, notes, idx, total)
            marker = "      "
        print(f"  + slide {idx} {marker} {data['title']}")

    prs.save(OUT_PPTX)
    print(f"\n[OK] wrote {OUT_PPTX}")
    print(f"     open: open '{OUT_PPTX}'")


if __name__ == "__main__":
    main()
